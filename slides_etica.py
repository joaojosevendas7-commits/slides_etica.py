from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Criar apresentação
prs = Presentation()

# Cores padrão
bg_color = RGBColor(30, 42, 56)  # azul escuro
text_color = RGBColor(255, 255, 255)  # branco
title_color = RGBColor(201, 162, 39)  # dourado

slides_content = [

("ÉTICA CRISTÃ",
"Definição de Ética, Moral e Ética Cristã\nProfessor João N."),

("O QUE É MORAL",
"Moral é o conjunto de valores, normas e costumes que orientam o comportamento humano.\n\nRomanos 2:15"),

("CARACTERÍSTICAS DA MORAL",
"- Prática\n- Vivida\n- Cultural\n- Social\nToda sociedade possui moral."),

("O QUE É ÉTICA",
"Ethos = hábito ou caráter.\nÉtica é o estudo dos princípios que orientam a conduta humana.\nProvérbios 4:23"),

("MORAL x ÉTICA",
"Moral diz o que fazemos.\nÉtica pergunta por que fazemos.\n1 Tessalonicenses 5:21"),

("SISTEMAS ÉTICOS",
"Existem três grandes sistemas:\n- Humanístico\n- Naturalista\n- Religioso"),

("ÉTICA HUMANÍSTICA",
"Baseada no homem como medida do certo.\nExemplos:\nHedonismo\nUtilitarismo\nExistencialismo"),

("ÉTICA NATURALISTA",
"Baseada na natureza.\nProblema: nem tudo que é natural é moral.\nRomanos 8:20"),

("ÉTICA RELIGIOSA",
"Baseada na revelação divina.\nDeus define o que é certo."),

("ÉTICA CRISTÃ",
"Estudo da conduta humana à luz da revelação de Deus.\nSalmo 119:105"),

("FUNDAMENTO DA MORAL CRISTÃ",
"O certo nasce do caráter de Deus.\nLevítico 19:2"),

("ÉTICA DO MUNDO x CRISTÃ",
"Mundo:\nRelativa\nMutável\n\nCristã:\nAbsoluta\nImutável"),

("MORAL DO MUNDO x MORAL CRISTÃ",
"Mundo:\nCultural\nRelativa\n\nCristã:\nBíblica\nPermanente"),

("COMO PENSAR ETICAMENTE",
"1. O que a Bíblia ensina?\n2. Qual princípio está envolvido?\n3. Isso glorifica a Deus?\n4. Isso ama o próximo?\n1 Coríntios 10:31"),

("FUNDAMENTOS DA ÉTICA CRISTÃ",
"1. Deus existe\n2. Deus se revela\n3. O homem é moral\n4. O homem é responsável\n5. Existe lei moral\n6. A vida tem propósito"),

("CONCLUSÃO",
"A ética cristã não é apenas saber o que é certo.\nÉ viver corretamente diante de Deus.")

]

for title_text, content_text in slides_content:

    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    title = slide.shapes.title
    content = slide.placeholders[1]

    title.text = title_text
    content.text = content_text

    # Estilo do título
    title_paragraph = title.text_frame.paragraphs[0]
    title_paragraph.font.size = Pt(44)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = title_color

    # Estilo do conteúdo
    content_paragraph = content.text_frame.paragraphs[0]
    content_paragraph.font.size = Pt(28)
    content_paragraph.font.color.rgb = text_color

# Salvar arquivo
prs.save("Aula_Etica_Crista.pptx")

print("Slides criados com sucesso!")
